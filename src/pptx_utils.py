from PIL import Image
from pptx.util import Mm

SLIDE_WIDTH, SLIDE_HEIGHT = 12193200, 6858000
# スライド中心のX、Y座標（左上が原点）
IMG_CENTER_X, IMG_CENTER_Y = SLIDE_WIDTH / 2, SLIDE_HEIGHT / 2
# スライドのアスペクト比
SLIDE_ASPECT_RATIO = SLIDE_WIDTH / SLIDE_HEIGHT


def add_slide(prs):
    # 白紙スライドの追加(ID=6は白紙スライド)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    return slide


def add_text(slide, txt, left=Mm(20), top=Mm(10), width=Mm(10), height=Mm(10)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = txt


def add_picture(slide, img_file):
    # 画像サイズを取得してアスペクト比を得る
    im = Image.open(img_file)
    im_width, im_height = im.size
    aspect_ratio = im_width / im_height

    # スライドと画像のアスペクト比に応じて処理を分岐
    # 画像のほうが横長だったら横めいっぱいに広げる
    if aspect_ratio > SLIDE_ASPECT_RATIO:
        img_display_width = SLIDE_WIDTH
        img_display_height = img_display_width / aspect_ratio
    else:  # 画像のほうが縦長だったら縦めいっぱいに広げる
        img_display_height = SLIDE_HEIGHT
        img_display_width = img_display_height * aspect_ratio
    # センタリングする場合の画像の左上座標を計算
    left = IMG_CENTER_X - img_display_width / 2
    top = IMG_CENTER_Y - img_display_height / 2

    # 画像をスライドに追加
    if aspect_ratio > SLIDE_ASPECT_RATIO:
        slide.shapes.add_picture(img_file, left, top, width=img_display_width)
    else:
        slide.shapes.add_picture(img_file, left, top, height=img_display_height)

    return slide
